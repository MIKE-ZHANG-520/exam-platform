from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
PRIMARY = RGBColor(0x16, 0x77, 0xFF)      # #1677ff 科技蓝
PRIMARY_DARK = RGBColor(0x09, 0x58, 0xD9)  # #0958d9
ACCENT = RGBColor(0xF2, 0x6E, 0x22)        # #F26E22 提示橘
SUCCESS = RGBColor(0x12, 0xA1, 0x50)       # 通过绿
DANGER = RGBColor(0xDC, 0x26, 0x26)        # 危险红
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)          # 标题灰
BODY = RGBColor(0x4B, 0x55, 0x63)          # 正文灰
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)      # 浅灰背景
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

SCREENSHOT_DIR = "/workspace/projects/public/screenshots"
OUTPUT_PATH = "/workspace/projects/public/智慧培训考试平台介绍.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_background(slide, color=LIGHT_BG):
    """Add solid background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """Add colored shape"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """Add text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=14, color=BODY, line_spacing=1.5):
    """Add multiline text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(font_size * 0.5)
    return txBox

def add_image(slide, image_path, left, top, width=None, height=None):
    """Add image with error handling"""
    if os.path.exists(image_path):
        if width and height:
            slide.shapes.add_picture(image_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(image_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            slide.shapes.add_picture(image_path, left, top)
        return True
    return False

# ============================================================
# SLIDE 1: Cover Page
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide1, WHITE)

# Left blue gradient bar
add_shape(slide1, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

# Main title area - left side
add_shape(slide1, Inches(0.8), Inches(1.5), Inches(6), Inches(0.08), PRIMARY)

add_text_box(slide1, Inches(0.8), Inches(1.8), Inches(6), Inches(1.2),
             "智慧培训考试平台", font_size=44, color=PRIMARY, bold=True)

add_text_box(slide1, Inches(0.8), Inches(3.0), Inches(6), Inches(0.8),
             "Smart Training & Examination Platform", font_size=20, color=BODY)

add_text_box(slide1, Inches(0.8), Inches(3.8), Inches(6), Inches(1.5),
             "AI 驱动 · 全流程数字化 · 安全合规闭环", font_size=16, color=ACCENT)

# Feature highlights
features = [
    "📄 材料智能解析 → AI 自动生成提纲与题库",
    "📱 扫码即考 → 手机 H5 答题，零安装门槛",
    "🛡️ 防作弊机制 → 倒计时 + 切屏检测 + 次数管控",
    "📊 数据看板 → 通过率 / 班组排名 / 趋势分析",
    " 工人档案 → 三级教育 / 入场交底 / 审核二维码",
]
add_multiline_text(slide1, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
                   features, font_size=13, color=BODY)

# Right side - login screenshot
add_image(slide1, f"{SCREENSHOT_DIR}/01_login.png",
          Inches(7.2), Inches(0.8), width=Inches(5.5))

# Bottom info
add_text_box(slide1, Inches(0.8), Inches(6.8), Inches(4), Inches(0.4),
             "企业级安全生产培训数字化解决方案", font_size=11, color=BODY)

print("✓ Slide 1: Cover")

# ============================================================
# SLIDE 2: AI-Powered Content Generation
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)
add_shape(slide2, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

# Section header
add_shape(slide2, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.5), PRIMARY)
add_text_box(slide2, Inches(0.9), Inches(0.45), Inches(5), Inches(0.6),
             "01  AI 驱动的内容生产引擎", font_size=28, color=PRIMARY, bold=True)
add_text_box(slide2, Inches(0.9), Inches(1.0), Inches(8), Inches(0.5),
             "上传一份 PDF/Word/Excel 培训材料，系统自动完成解析、提纲生成、题库构建——将传统数天的备课工作压缩至分钟级", font_size=13, color=BODY)

# Left: Materials page screenshot
add_image(slide2, f"{SCREENSHOT_DIR}/03_materials.png",
          Inches(0.6), Inches(1.6), width=Inches(6.2))

# Right: Feature cards
card_data = [
    ("智能解析", "支持 PDF / Word / Excel\n多格式文件自动提取文本\n外部解析服务兜底", PRIMARY),
    ("双版本提纲", "工人版：故事化场景教学\n培训师版：互动化授课框架\n适配不同受众认知习惯", SUCCESS),
    ("AI 出题", "按难度分级（基础/中等）\n单选 / 多选 / 判断题\n人工审核后可入库", ACCENT),
]

for i, (title, desc, color) in enumerate(card_data):
    y = Inches(1.6) + Inches(i * 1.85)
    # Card background
    card = add_shape(slide2, Inches(7.2), y, Inches(5.5), Inches(1.6), CARD_BG)
    card.shadow.inherit = False
    # Left color bar
    add_shape(slide2, Inches(7.2), y, Inches(0.08), Inches(1.6), color)
    # Title
    add_text_box(slide2, Inches(7.5), y + Inches(0.15), Inches(2), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    # Description
    add_multiline_text(slide2, Inches(7.5), y + Inches(0.55), Inches(5), Inches(1),
                       desc.split('\n'), font_size=11, color=BODY)

print("✓ Slide 2: AI Content")

# ============================================================
# SLIDE 3: Exam System & Anti-Cheating
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)
add_shape(slide3, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

add_shape(slide3, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.5), PRIMARY)
add_text_box(slide3, Inches(0.9), Inches(0.45), Inches(5), Inches(0.6),
             "02  闭环考试体系与防作弊机制", font_size=28, color=PRIMARY, bold=True)
add_text_box(slide3, Inches(0.9), Inches(1.0), Inches(8), Inches(0.5),
             "从试卷创建到扫码考试、自动评分、成绩归档，全流程数字化；内置多重防作弊手段保障考试公平性", font_size=13, color=BODY)

# Left column: Exam management screenshot
add_image(slide3, f"{SCREENSHOT_DIR}/05_exams.png",
          Inches(0.6), Inches(1.6), width=Inches(5.8))

# Right column: Anti-cheating features
anti_cheat = [
    ("⏱ 倒计时控场", "限时答题，超时自动交卷\n杜绝拖延与场外求助"),
    (" 切屏检测", "实时监测页面切换行为\n记录切屏次数，超限预警"),
    ("🔄 12小时窗口制", "同一试卷12小时内限考2次\n超时自动重置，兼顾管控与灵活"),
    ("📋 试卷快照", "每次考试随机抽题生成快照\n答案绑定快照，防止篡改"),
    ("🏷 A/B卷机制", "支持多套试卷配置\n不同场次使用不同试卷"),
]

for i, (title, desc) in enumerate(anti_cheat):
    y = Inches(1.6) + Inches(i * 1.1)
    add_text_box(slide3, Inches(6.8), y, Inches(3), Inches(0.35),
                 title, font_size=14, color=DARK, bold=True)
    add_multiline_text(slide3, Inches(6.8), y + Inches(0.35), Inches(5.5), Inches(0.7),
                       desc.split('\n'), font_size=11, color=BODY)

# Bottom: Records screenshot
add_image(slide3, f"{SCREENSHOT_DIR}/06_records.png",
          Inches(6.8), Inches(5.2), width=Inches(5.8))

print("✓ Slide 3: Exam System")

# ============================================================
# SLIDE 4: Worker Safety Management
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)
add_shape(slide4, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

add_shape(slide4, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.5), PRIMARY)
add_text_box(slide4, Inches(0.9), Inches(0.45), Inches(5), Inches(0.6),
             "03  工人安全档案全生命周期管理", font_size=28, color=PRIMARY, bold=True)
add_text_box(slide4, Inches(0.9), Inches(1.0), Inches(8), Inches(0.5),
             "从花名册导入到三级安全教育、入场交底、专项培训、档案审核，构建完整的工人安全准入闭环", font_size=13, color=BODY)

# Left: Safety management page
add_image(slide4, f"{SCREENSHOT_DIR}/07_safety.png",
          Inches(0.6), Inches(1.6), width=Inches(6.2))

# Right: Lifecycle stages
stages = [
    ("① 花名册导入", "Excel/CSV 批量导入\nAI 智能识别表头与列名\n自动匹配班组与工种", PRIMARY),
    (" 三级安全教育", "公司级 → 项目级 → 班组级\n逐级培训记录留痕\n培训时长与内容可追溯", SUCCESS),
    ("③ 入场安全交底", "交底内容分类上传\n综合资料 / 教育记录 / 交底文件\n支持 PDF/图片多格式", ACCENT),
    ("④ 审核与二维码", "管理员审核档案完整性\n通过后生成个人二维码\n扫码即可查看完整档案", DANGER),
]

for i, (title, desc, color) in enumerate(stages):
    y = Inches(1.6) + Inches(i * 1.4)
    card = add_shape(slide4, Inches(7.2), y, Inches(5.5), Inches(1.2), CARD_BG)
    add_shape(slide4, Inches(7.2), y, Inches(0.08), Inches(1.2), color)
    add_text_box(slide4, Inches(7.5), y + Inches(0.1), Inches(2.5), Inches(0.35),
                 title, font_size=15, color=color, bold=True)
    add_multiline_text(slide4, Inches(7.5), y + Inches(0.45), Inches(5), Inches(0.7),
                       desc.split('\n'), font_size=11, color=BODY)

# Bottom: Operation logs
add_image(slide4, f"{SCREENSHOT_DIR}/08_logs.png",
          Inches(0.6), Inches(5.3), width=Inches(6.2))

add_text_box(slide4, Inches(7.2), Inches(5.3), Inches(5.5), Inches(0.4),
             "📝 操作日志：关键操作全程留痕，满足审计合规要求", font_size=12, color=BODY)

print("✓ Slide 4: Worker Safety")

# ============================================================
# SLIDE 5: Data Insights & Management Value
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)
add_shape(slide5, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

add_shape(slide5, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.5), PRIMARY)
add_text_box(slide5, Inches(0.9), Inches(0.45), Inches(5), Inches(0.6),
             "04  数据驱动的安全管理决策", font_size=28, color=PRIMARY, bold=True)
add_text_box(slide5, Inches(0.9), Inches(1.0), Inches(8), Inches(0.5),
             "从经验管理走向数据管理，让安全培训效果可量化、可对比、可改进", font_size=13, color=BODY)

# Dashboard screenshot - full width
add_image(slide5, f"{SCREENSHOT_DIR}/02_dashboard.png",
          Inches(0.6), Inches(1.6), width=Inches(12))

# Bottom: Value propositions
values = [
    ("📈 实时看板", "参考人数 / 通过率 / 平均分数\n班组排名对比 / 日趋势分析\n一眼掌握培训全局态势"),
    ("🎯 精准干预", "低分人员自动标记待补考\n按班组维度定位薄弱环节\n针对性强化培训，资源精准投放"),
    ("📋 合规留痕", "考试记录 / 培训记录 / 操作日志\n全部电子化存档，随时调阅\n满足安监检查与审计要求"),
    ("⚡ 效率提升", "AI 出题节省 90% 备课时间\n扫码考试零组织成本\n数据自动汇总，告别手工统计"),
]

for i, (title, desc) in enumerate(values):
    x = Inches(0.6) + Inches(i * 3.15)
    y = Inches(5.2)
    add_text_box(slide5, x, y, Inches(3), Inches(0.35),
                 title, font_size=14, color=PRIMARY, bold=True)
    add_multiline_text(slide5, x, y + Inches(0.35), Inches(3), Inches(1.2),
                       desc.split('\n'), font_size=11, color=BODY)

print("✓ Slide 5: Data Insights")

# Save
prs.save(OUTPUT_PATH)
print(f"\n✅ PPT saved to: {OUTPUT_PATH}")
print(f"   File size: {os.path.getsize(OUTPUT_PATH) / 1024:.0f} KB")
